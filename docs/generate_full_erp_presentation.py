"""
Generate complete Shree Hari Jewels ERP + Website guide with screenshots.
Run: python docs/generate_full_erp_presentation.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

DOCS = Path(__file__).parent
ERP_SHOTS = DOCS / "screenshots" / "erp"
SHOP_SHOTS = DOCS / "screenshots"
OUTPUT = DOCS / "Shri-Hari-Jewels-Complete-ERP-Guide.pptx"

GOLD = RGBColor(0xB8, 0x86, 0x0B)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_BG = RGBColor(0xF8, 0xF5, 0xEF)

STORE_URL = "https://shri-hari-jewels.vercel.app/shop/shree-hari-jewels"
ERP_URL = "https://shri-hari-jewels.vercel.app"


def set_slide_bg(slide, color=LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, prs, title):
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK
    header.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.16), Inches(12.5), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.3))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.3), Inches(1.8))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = GOLD
        p2.alignment = PP_ALIGN.CENTER


def add_section_slide(prs, part, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(3.0), Inches(0.1), Inches(1.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    if part:
        pt = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(0.5))
        pp = pt.text_frame.paragraphs[0]
        pp.text = part
        pp.font.size = Pt(14)
        pp.font.color.rgb = GOLD
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12), Inches(0.8))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(15)
        p2.font.color.rgb = GOLD


def add_steps_slide(prs, title, steps, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, prs, title)
    box = slide.shapes.add_textbox(Inches(0.45), Inches(1.05), Inches(12.4), Inches(5.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, step in enumerate(steps):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = f"Step {i + 1}:  {step}"
        para.font.size = Pt(14)
        para.font.color.rgb = DARK
        para.space_after = Pt(8)
    if note:
        nb = slide.shapes.add_textbox(Inches(0.45), Inches(6.55), Inches(12), Inches(0.45))
        np = nb.text_frame.paragraphs[0]
        np.text = note
        np.font.size = Pt(11)
        np.font.italic = True
        np.font.color.rgb = GOLD


def find_image(name):
    for folder in (ERP_SHOTS, SHOP_SHOTS):
        p = folder / f"{name}.png"
        if p.exists():
            return p
    # fallback: old naming
    old = SHOP_SHOTS / name
    if old.exists():
        return old
    return None


def add_screenshot_slide(prs, title, image_name, steps=None, caption="", erp=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, prs, title)

    img_path = find_image(image_name)
    if img_path and img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.3), Inches(1.0), width=Inches(8.0))
        right_x = Inches(8.5)
        width = Inches(4.5)
    else:
        right_x = Inches(0.45)
        width = Inches(12.2)
        warn = slide.shapes.add_textbox(Inches(0.45), Inches(1.2), Inches(8), Inches(0.4))
        warn.text_frame.paragraphs[0].text = f"[Screenshot: {image_name}.png]"
        warn.text_frame.paragraphs[0].font.color.rgb = GRAY

    if steps:
        box = slide.shapes.add_textbox(right_x, Inches(1.05), width, Inches(5.7))
        tf = box.text_frame
        tf.word_wrap = True
        h = tf.paragraphs[0]
        h.text = "Steps:"
        h.font.bold = True
        h.font.size = Pt(12)
        h.font.color.rgb = GOLD
        h.space_after = Pt(5)
        for step in steps:
            para = tf.add_paragraph()
            para.text = f"• {step}"
            para.font.size = Pt(10)
            para.font.color.rgb = GRAY
            para.space_after = Pt(4)

    if caption:
        cap = slide.shapes.add_textbox(Inches(0.3), Inches(6.8), Inches(12.5), Inches(0.4))
        cp = cap.text_frame.paragraphs[0]
        cp.text = caption
        cp.font.size = Pt(9)
        cp.font.italic = True
        cp.font.color.rgb = GRAY


def add_flow_slide(prs, title, columns):
    """columns: list of (heading, [items])"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, prs, title)
    n = len(columns)
    col_w = 12.3 / n
    for i, (heading, items) in enumerate(columns):
        x = Inches(0.35 + i * col_w)
        w = Inches(col_w - 0.12)
        hdr = slide.shapes.add_shape(1, x, Inches(1.05), w, Inches(0.4))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = GOLD
        hdr.line.fill.background()
        ht = slide.shapes.add_textbox(x, Inches(1.08), w, Inches(0.35))
        hp = ht.text_frame.paragraphs[0]
        hp.text = heading
        hp.font.bold = True
        hp.font.size = Pt(10)
        hp.font.color.rgb = WHITE
        hp.alignment = PP_ALIGN.CENTER
        body = slide.shapes.add_textbox(x, Inches(1.55), w, Inches(5.3))
        btf = body.text_frame
        btf.word_wrap = True
        for j, item in enumerate(items):
            para = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
            para.text = f"{j + 1}. {item}"
            para.font.size = Pt(9)
            para.font.color.rgb = GRAY
            para.space_after = Pt(3)


# ── Module definitions: (section, slides[]) ──
# Each slide: type = steps|screenshot|flow, title, content...

MODULES = []  # filled in build_presentation


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # TITLE
    add_title_slide(
        prs,
        "Shree Hari Jewels — Complete ERP Guide",
        "Every Module · Every Step · With Screenshots\n"
        f"ERP: {ERP_URL}  |  Store: {STORE_URL}",
    )

    # TOC
    add_steps_slide(
        prs,
        "What's Inside This Guide",
        [
            "Part 1 — Getting Started: Login, Dashboard, Settings, Branches",
            "Part 2 — Central Stock: Add products, add units, entry verification",
            "Part 3 — Raw Materials: Metal lots, certified stones, bulk stone stock",
            "Part 4 — Production: Motifs, Designs (6-stage builder), Production Runs (10 stages), Karigar settlements",
            "Part 5 — Stock Transfer: Scan & send, sent list, proforma, incoming receive",
            "Part 6 — Sales & CRM: POS, custom orders, customers, invoices",
            "Part 7 — Reports: GST, stock valuation, ageing, analytics",
            "Part 8 — Online Store: Website setup, customer journey, web order fulfillment",
        ],
    )

    # ═══ PART 1: GETTING STARTED ═══
    add_section_slide(prs, "Part 1", "Getting Started", "Login · Dashboard · Settings · Branches")

    add_steps_slide(
        prs,
        "How to Log In",
        [
            f"Open {ERP_URL}/login in your browser",
            "Enter your User ID or email (e.g. admin@shreehari.com)",
            "Enter your password → click Sign In",
            "If 2FA is enabled, enter the 6-digit code from your authenticator app",
            "You land on the Dashboard — your home screen",
        ],
        note="Each staff member has a role (Admin, Store, Production Manager, etc.) that controls which menus they see.",
    )

    add_screenshot_slide(
        prs, "Dashboard — Your Home Screen", "erp-01-dashboard",
        ["See overview stats: stock, sales, pending orders", "Market rate banner shows live gold/silver prices",
         "Sidebar gives access to all modules", "Switch branch from top bar if you have multiple locations"],
        "ERP → Dashboard",
    )

    add_screenshot_slide(
        prs, "Settings — Business Configuration", "erp-33-settings",
        ["Set business name, GST number, PAN, bank details", "Configure making charge % for gold and silver",
         "Set discount approval threshold (discounts above this need admin approval)",
         "Manage users: create staff, assign roles and branches"],
        "ERP → System → Settings",
    )

    add_screenshot_slide(
        prs, "Branches — Store Locations", "erp-32-branches",
        ["View all branches: head office + store locations", "Add new branch with name, address, GST",
         "Assign staff users to specific branches", "Stock and sales are scoped by branch"],
        "ERP → System → Branches",
    )

    # ═══ PART 2: CENTRAL STOCK ═══
    add_section_slide(prs, "Part 2", "Central Stock (Finished Goods)",
                      "Every piece tracked by unique item code")

    add_steps_slide(
        prs,
        "Stock Entry — Full Workflow",
        [
            "Create a Product SKU (metal, purity, weight, making charges) → /inventory/new",
            "Add physical units to that SKU → each unit gets a unique item code",
            "Units go to Pending Verification status → appear in Entry Verification",
            "Verifier sets list price per unit using live market rates",
            "After verification → units become Available and can be sold or transferred",
        ],
        note="List price = market rate × weight + making charges (auto-calculated).",
    )

    add_screenshot_slide(
        prs, "All Stock — View Central Inventory", "erp-02-inventory",
        ["See all products with stock counts", "Filter by metal, category, branch, status",
         "Click item code to see full history (sales, transfers, status changes)",
         "Click Add Stock to create new product SKU"],
        "ERP → Inventory → All stock",
    )

    add_screenshot_slide(
        prs, "Add Stock — Create New Product SKU", "erp-03-add-stock",
        ["Enter product name, SKU code, category", "Select metal type and purity (22K Gold, 925 Silver, etc.)",
         "Enter gross weight and making charge %", "Upload product images", "Save → product created (0 units initially)"],
        "ERP → Inventory → Add Stock",
    )

    add_steps_slide(
        prs,
        "Add Units to a Product",
        [
            "From All Stock, open a product → click Add Units",
            "Scan or type item codes for each physical piece (or auto-generate)",
            "Each unit creates an Entry Voucher line (status: Pending Verification)",
            "All units for this batch share one Entry Voucher",
            "Go to Entry Verification to set prices and approve",
        ],
    )

    add_screenshot_slide(
        prs, "Entry Verification — Approve New Stock", "erp-04-entry-verification",
        ["Lists all pending entry vouchers awaiting verification", "Open a voucher → see all units with weights",
         "System suggests list price from market rates — adjust if needed",
         "Click Verify → all units become Available", "Product stock count updates automatically"],
        "ERP → Inventory → Entry verification",
    )

    # ═══ PART 3: RAW MATERIALS ═══
    add_section_slide(prs, "Part 3", "Raw Materials",
                      "Metal lots · Certified stones · Bulk stone stock")

    add_steps_slide(
        prs,
        "Raw Materials — What They Are Used For",
        [
            "Metal Lots: raw gold/silver/platinum purchased from vendors — used in production runs",
            "Certified Stones: GIA/IGI certified diamonds/gems with certificate numbers",
            "Bulk Stone Stock: uncertified stones tracked by pieces or carat weight",
            "When a production run starts, metal is deducted automatically from metal lots",
            "Stones are issued to karigars during the Stone Setting production stage",
        ],
    )

    add_screenshot_slide(
        prs, "Raw Inventory — Metal & Stones", "erp-05-raw-inventory",
        ["Metal Lots tab: add gold/silver lots with purity, weight, purchase rate, vendor",
         "Certified Stones tab: add stones with certificate, carat, color, clarity",
         "Bulk Stone Stock tab: track stones by type, pieces, weight, rate basis",
         "Market rates banner shows current gold/silver prices used everywhere"],
        "ERP → Inventory → Raw materials",
    )

    add_steps_slide(
        prs,
        "Add a Metal Lot",
        [
            "Click Add Metal Lot",
            "Select metal type (Gold / Silver / Platinum) and purity (22K, 18K, 925, etc.)",
            "Enter weight in grams, purchase rate per gram, vendor name",
            "Select storage location (branch)",
            "Save → lot available for production run deduction",
        ],
    )

    # ═══ PART 4: PRODUCTION ═══
    add_section_slide(prs, "Part 4", "Production",
                      "Motifs → Designs → Production Runs → Karigar Settlements")

    add_steps_slide(
        prs,
        "Production Pipeline Overview",
        [
            "Step A: Build Motif Library — reusable components (pendant motif, earring hook, etc.)",
            "Step B: Create Design — 6-stage builder (SKU → CAD → Motifs → Mold → Photo → Complete)",
            "Step C: Submit design for approval (Approved required before production)",
            "Step D: Create Production Run — select design + number of sets",
            "Step E: Progress through 10 workshop stages with karigars",
            "Step F: On completion → finished goods auto-created in central stock",
            "Step G: Karigar settlements — pay wastage cost + making wages",
        ],
    )

    add_screenshot_slide(
        prs, "Motif Library", "erp-17-motifs",
        ["Reusable jewellery components with metal, weight, stones, making cost",
         "Price auto-calculates from market rates", "Import motifs from Excel in bulk",
         "Used as building blocks when creating design BOMs"],
        "ERP → Production → Motifs",
    )

    add_steps_slide(
        prs,
        "Create a New Design — 6 Stages",
        [
            "Stage 1 SKU: Define design code, category, target metal and weight",
            "Stage 2 CAD: Upload CAD files, mark design as CAD-ready",
            "Stage 3 Motifs: Add motif elements to BOM — qty, weight, unit value per element",
            "Stage 4 Mold Making: Track mold creation",
            "Stage 5 Photo: Upload product photos",
            "Stage 6 Complete: Submit for approval → Admin approves or rejects",
        ],
    )

    add_screenshot_slide(
        prs, "Designs List & New Design", "erp-15-designs",
        ["See all designs with approval status (Draft / Pending / Approved / Rejected)",
         "Click New Design to start the 6-stage builder",
         "Open a design → walk through each builder stage",
         "BOM shows total metal and stone requirements per piece"],
        "ERP → Production → Designs",
    )

    add_screenshot_slide(
        prs, "New Design Form", "erp-16-new-design",
        ["Enter design code and name", "Select category and metal type",
         "Set target weight — used to calculate BOM requirements",
         "After saving, open Design Builder to complete all 6 stages"],
        "ERP → Production → Designs → New Design",
    )

    add_steps_slide(
        prs,
        "Create a Production Run",
        [
            "Go to Production Runs → New Production Run",
            "Select an Approved design",
            "Enter number of sets (pieces) to manufacture",
            "System previews metal & stone requirements (BOM × sets)",
            "On create: metal deducted immediately from raw inventory",
            "Run appears on Production Board with all items to track",
        ],
    )

    add_screenshot_slide(
        prs, "Production Runs List", "erp-19-production-runs",
        ["All active and completed production runs", "See design, sets count, current stage, status",
         "Click a run to open stage worksheets", "New Production Run button top-right"],
        "ERP → Production → Production runs",
    )

    add_screenshot_slide(
        prs, "New Production Run", "erp-20-new-production-run",
        ["Select approved design from dropdown", "Enter number of sets to produce",
         "Preview shows total metal and stone needed", "Confirm → run created, metal reserved/deducted"],
        "ERP → Production → Production runs → New",
    )

    add_steps_slide(
        prs,
        "10 Production Stages (Workshop Workflow)",
        [
            "1. Wax Pattern — create wax model from mold",
            "2. Casting — metal casting; issue metal to karigar",
            "3. Filing — shape and file the cast piece",
            "4. Soldering — join components, add findings",
            "5. Stone Setting — set stones; issue stone stock to karigar",
            "6. Polishing — final polish and finish",
            "7. Rhodium / Plating — optional plating step",
            "8. QC — quality check; pass or reject for rework",
            "9. Hallmarking — BIS hallmark if required",
            "10. Packaging — final pack → triggers finished goods creation in stock",
        ],
    )

    add_screenshot_slide(
        prs, "Production Board — Active Runs Dashboard", "erp-21-production-board",
        ["Visual board of all active production runs", "See which stage each run is at",
         "Click a run to open its stage worksheet", "Karigars update items at each stage"],
        "ERP → Production → Production board",
    )

    add_steps_slide(
        prs,
        "Karigar Settlements",
        [
            "After each production stage, record metal issued to and returned from karigar",
            "System calculates metal wastage (issued − returned − expected)",
            "Wastage cost = wastage weight × current metal rate",
            "Making wage = per-piece making charge × pieces completed",
            "Karigar Settlements page shows all pending settlements",
            "Mark as paid once settled with the karigar",
        ],
    )

    add_screenshot_slide(
        prs, "Karigar Settlements", "erp-22-karigar-settlements",
        ["Lists all karigar settlements by production run and stage",
         "Shows wastage weight, wastage cost, making wage payable",
         "Filter by karigar name or run", "Mark settlement as Paid when done"],
        "ERP → Production → Karigar settlements",
    )

    add_screenshot_slide(
        prs, "Work Orders", "erp-18-work-orders",
        ["Internal workshop tasks — repairs, custom jobs, modifications",
         "Create work order → assign to karigar", "Track status: Open → In Production → QC → Completed",
         "Can link to a customer order for custom jewellery jobs"],
        "ERP → Production → Work orders",
    )

    # ═══ PART 5: STOCK TRANSFER ═══
    add_section_slide(prs, "Part 5", "Stock Transfer",
                      "Move stock from head office to branch stores")

    add_steps_slide(
        prs,
        "Stock Transfer — Full Workflow",
        [
            "Head office: Scan & Send — scan item codes, select destination branch",
            "Choose document type: Wholesale GST Invoice / Delivery Challan / Stock Transfer Note",
            "Transfer created → units status changes to In Transit",
            "Fill shipping details and GST fields → invoice PDF generated",
            "Receiving branch: Incoming — scan and accept (full or partial) or reject items",
            "Accepted units become Available at the destination branch",
        ],
    )

    add_screenshot_slide(
        prs, "Scan & Send — Create a Transfer", "erp-06-scan-send",
        ["Scan or type item codes of units to transfer", "Select destination branch or customer branch",
         "Choose document type (GST Invoice / Challan / STN)", "Review items → Create Transfer",
         "Units immediately go to In Transit status"],
        "ERP → Stock transfer → Scan & send",
    )

    add_screenshot_slide(
        prs, "Sent Transfers — Track Outgoing Stock", "erp-07-sent-transfers",
        ["All transfers sent from your branch", "See status: Pending / Accepted / Rejected / Partial",
         "Click a transfer to add shipping details or download invoice PDF",
         "Proforma list shows transfers awaiting acceptance"],
        "ERP → Stock transfer → Sent",
    )

    add_screenshot_slide(
        prs, "Proforma List — Pending Transfer Invoices", "erp-08-proforma",
        ["Transfers sent but not yet accepted by receiving branch",
         "Fill courier/shipping details before sending goods",
         "Generate and download GST invoice or delivery challan PDF"],
        "ERP → Stock transfer → Proforma list",
    )

    add_screenshot_slide(
        prs, "Incoming Stock — Receive at Branch", "erp-09-incoming",
        ["Transfers sent TO your branch appear here", "Open a transfer → scan item codes to receive",
         "Accept all, accept partial, or reject damaged items",
         "Accepted units become Available at your branch"],
        "ERP → Stock transfer → Incoming",
    )

    # ═══ PART 6: SALES & CRM ═══
    add_section_slide(prs, "Part 6", "Sales & CRM",
                      "POS · Custom Orders · Customers · Invoices")

    add_steps_slide(
        prs,
        "In-Store Sale (POS) — Step by Step",
        [
            "Open Sales page → scan physical item code (barcode or manual entry)",
            "System shows live price breakdown: metal value + making charges",
            "Add more items to cart if needed",
            "Apply discount — if above threshold, admin approval required",
            "Select or create customer (search by mobile / email / GST)",
            "Choose payment: Cash / UPI (Razorpay QR) / Card",
            "Complete sale → unit marked Sold, GST invoice auto-generated",
        ],
    )

    add_screenshot_slide(
        prs, "Sales — Point of Sale", "erp-10-sales",
        ["Scan item code → item appears in cart with live price",
         "Apply discount, select customer, choose payment mode",
         "UPI: QR code generated → customer pays → webhook confirms",
         "Old gold exchange credit can be applied to reduce total"],
        "ERP → Sales → Sales (POS)",
    )

    add_steps_slide(
        prs,
        "Custom Customer Order — Step by Step",
        [
            "Go to Orders → New Order",
            "Select existing customer or create new one",
            "Enter order description, estimated total, due date, notes",
            "Order status progresses: Pending → Designing → Production → QC → Ready → Delivered",
            "Create Work Orders and/or Production Runs linked to this order",
            "Track payment: Unpaid / Partial / Paid",
        ],
    )

    add_screenshot_slide(
        prs, "Orders List", "erp-11-orders",
        ["All custom jewellery orders with status and due dates",
         "Filter by status, customer, date range", "Click New Order to create"],
        "ERP → Sales → Orders",
    )

    add_screenshot_slide(
        prs, "New Order Form", "erp-12-new-order",
        ["Select customer from dropdown or create inline",
         "Enter description of custom jewellery piece required",
         "Set estimated total and due date", "Save → order created in Pending status"],
        "ERP → Sales → Orders → New Order",
    )

    add_screenshot_slide(
        prs, "Customers — CRM", "erp-13-customers",
        ["All customers: individual buyers and B2B clients",
         "Add customer with billing address, GST/PAN, bank details",
         "Customer branches for B2B clients with multiple offices",
         "Customer record shared across POS, orders, and web checkout"],
        "ERP → Sales → Customers",
    )

    add_screenshot_slide(
        prs, "Invoices — GST Billing", "erp-14-invoices",
        ["All GST invoices from POS sales and wholesale transfers",
         "Download invoice PDF with CGST/SGST or IGST split",
         "Generate e-invoice IRN for GST portal compliance",
         "Filter by date, customer, branch, payment status"],
        "ERP → Sales → Invoices",
    )

    # ═══ PART 7: REPORTS ═══
    add_section_slide(prs, "Part 7", "Reports & Analytics")

    add_screenshot_slide(
        prs, "Sales Analytics", "erp-23-sales-analytics",
        ["Revenue charts by day/week/month", "Top selling products and categories",
         "Sales by branch and staff member", "Export data for further analysis"],
        "ERP → Reports → Sales analytics",
    )

    add_screenshot_slide(
        prs, "GST Report", "erp-24-gst-report",
        ["Tax liability summary by period", "CGST, SGST, IGST breakdown",
         "Export for GST filing", "Filter by date range and branch"],
        "ERP → Reports → GST report",
    )

    add_screenshot_slide(
        prs, "Stock Valuation Report", "erp-25-stock-valuation",
        ["Total inventory value at current market rates",
         "Breakdown by metal type and category", "Useful for balance sheet and insurance"],
        "ERP → Reports → Stock valuation",
    )

    add_screenshot_slide(
        prs, "Ageing Stock Report", "erp-26-ageing-stock",
        ["Items unsold beyond a threshold number of days",
         "Identify slow-moving inventory", "Plan promotions or transfers for stale stock"],
        "ERP → Reports → Ageing stock",
    )

    add_screenshot_slide(
        prs, "Staff Performance Report", "erp-27-staff-performance",
        ["Sales productivity by staff member", "Compare performance across branches",
         "Use for incentives and training decisions"],
        "ERP → Reports → Staff performance",
    )

    # ═══ PART 8: ONLINE STORE ═══
    add_section_slide(prs, "Part 8", "Online Store",
                      "Website setup · Customer journey · Web order fulfillment")

    add_flow_slide(
        prs,
        "Complete Business Flow — Design to Sale",
        [
            ("PRODUCTION", ["Create motifs", "Build & approve design", "Start production run",
                            "10 workshop stages", "Finished goods → stock"]),
            ("INVENTORY", ["Entry verification", "Set list prices", "Available units",
                           "Transfer to branches", "Publish to website"]),
            ("SALES", ["In-store POS scan", "Web order checkout", "GST invoice",
                       "E-invoice IRN", "Reports & analytics"]),
        ],
    )

    add_steps_slide(
        prs,
        "Set Up Your Online Store — 5 Steps",
        [
            "ERP → Online Store → Store Settings → Enable online store",
            "Set branding: logo, hero banner, colours, contact info, WhatsApp",
            "Publish Products → toggle inventory items to Published",
            "Create Collections to group products (optional)",
            f"Share store URL: {STORE_URL}",
        ],
    )

    add_screenshot_slide(
        prs, "Online Store Dashboard", "erp-28-storefront",
        ["Stats: published products, collections, pending web orders",
         "View Store button opens live website", "Quick links to all store management pages"],
        "ERP → Online Store → Store dashboard",
    )

    add_screenshot_slide(
        prs, "Store Settings", "erp-29-store-settings",
        ["Enable/disable store, set branding and theme colours",
         "Hero title, subtitle, about text, logo and banner URLs",
         "Contact phone, WhatsApp, address, return policy"],
        "ERP → Online Store → Store settings",
    )

    add_screenshot_slide(
        prs, "Publish Products to Website", "erp-30-publish-products",
        ["Toggle products from inventory to Published",
         "Only Published products with stock appear on website",
         "Prices sync from live market rates automatically"],
        "ERP → Online Store → Publish products",
    )

    add_steps_slide(
        prs,
        "Customer Website Journey — 7 Steps",
        [
            f"Customer opens {STORE_URL}",
            "Browses homepage, product catalog, or collections",
            "Opens product → sees specs, live price, stock status",
            "Adds to cart (saved in browser — no account needed)",
            "Proceeds to checkout → fills guest form (name, mobile, address)",
            "Clicks Place Order → order number shown on confirmation page",
            "Store contacts customer for payment and delivery arrangement",
        ],
    )

    add_screenshot_slide(
        prs, "Website — Homepage", "shop-01-home",
        ["Hero banner, featured products, brand story",
         "Navigation: Home, Shop, Collections, About", "Cart badge shows item count"],
        f"Customer → {STORE_URL}",
    )

    add_screenshot_slide(
        prs, "Website — Product Catalog", "shop-02-products",
        ["All published products in a grid", "Search, category filter, sort by price/name",
         "Live prices from market rates"],
        "Customer → Shop → All Products",
    )

    add_screenshot_slide(
        prs, "Website — Checkout", "shop-04-checkout",
        ["Guest form: name, mobile, address, city, state, pincode",
         "Order summary on right with total",
         "Payment arranged with store after order — no online payment gateway"],
        "Customer → Checkout",
    )

    add_steps_slide(
        prs,
        "Fulfill a Web Order — Step by Step",
        [
            "New order appears in ERP → Online Store → Web Orders",
            "Call/message customer to confirm order and arrange payment",
            "Update Order Status: Pending → Confirmed → Processing → Shipped → Delivered",
            "Update Payment Status: Unpaid → Paid (after bank transfer / COD)",
            "Inventory units move from Reserved → Sold when order is fulfilled",
            "Customer record auto-created in CRM from checkout mobile number",
        ],
    )

    add_screenshot_slide(
        prs, "Web Orders — Manage Online Orders", "erp-31-web-orders",
        ["All web orders with customer name, items, total, date",
         "Update order status and payment status dropdowns",
         "Linked ERP Order created automatically for fulfillment tracking"],
        "ERP → Online Store → Web orders",
    )

    # END
    add_title_slide(
        prs,
        "End of Guide",
        f"ERP: {ERP_URL}\nStore: {STORE_URL}\n\n"
        "Every module · Every step · Shree Hari Jewels",
    )

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
