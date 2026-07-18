"""
Generate Shree Hari Jewels Online Store + ERP workflow presentation.
Focused on single-company website operation with screenshots.

Run: python docs/generate_erp_presentation.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

DOCS = Path(__file__).parent
SCREENSHOTS = DOCS / "screenshots"
OUTPUT = DOCS / "Shri-Hari-Jewels-Online-Store-Guide.pptx"

GOLD = RGBColor(0xB8, 0x86, 0x0B)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_BG = RGBColor(0xF8, 0xF5, 0xEF)

STORE_URL = "https://shri-hari-jewels.vercel.app/shop/shree-hari-jewels"
ERP_URL = "https://shri-hari-jewels.vercel.app"


def set_slide_bg(slide, color=LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, prs, title):
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.95))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK
    header.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(12), Inches(1.4))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.6), Inches(12), Inches(1.5))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(17)
        p2.font.color.rgb = GOLD
        p2.alignment = PP_ALIGN.CENTER


def add_section_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(3.1), Inches(0.1), Inches(1.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(2.9), Inches(12), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.55), Inches(4.0), Inches(12), Inches(0.8))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = GOLD


def add_steps_slide(prs, title, steps, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, prs, title)

    box = slide.shapes.add_textbox(Inches(0.55), Inches(1.15), Inches(12.2), Inches(5.9))
    tf = box.text_frame
    tf.word_wrap = True
    for i, step in enumerate(steps):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = f"Step {i + 1}:  {step}"
        para.font.size = Pt(15)
        para.font.color.rgb = DARK
        para.space_after = Pt(10)

    if note:
        nb = slide.shapes.add_textbox(Inches(0.55), Inches(6.5), Inches(12), Inches(0.5))
        np = nb.text_frame.paragraphs[0]
        np.text = note
        np.font.size = Pt(12)
        np.font.italic = True
        np.font.color.rgb = GOLD


def add_screenshot_slide(prs, title, image_name, steps=None, caption=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, prs, title)

    img_path = SCREENSHOTS / image_name
    if img_path.exists():
        slide.shapes.add_picture(
            str(img_path),
            Inches(0.35),
            Inches(1.05),
            width=Inches(8.2),
        )
        right_x = Inches(8.75)
        width = Inches(4.2)
    else:
        right_x = Inches(0.55)
        width = Inches(12.0)
        warn = slide.shapes.add_textbox(Inches(0.55), Inches(1.2), Inches(8), Inches(0.4))
        warn.text_frame.paragraphs[0].text = f"[Screenshot missing: {image_name}]"
        warn.text_frame.paragraphs[0].font.color.rgb = GRAY

    if steps:
        box = slide.shapes.add_textbox(right_x, Inches(1.1), width, Inches(5.8))
        tf = box.text_frame
        tf.word_wrap = True
        heading = tf.paragraphs[0]
        heading.text = "What happens:"
        heading.font.bold = True
        heading.font.size = Pt(13)
        heading.font.color.rgb = GOLD
        heading.space_after = Pt(6)
        for step in steps:
            para = tf.add_paragraph()
            para.text = f"• {step}"
            para.font.size = Pt(11)
            para.font.color.rgb = GRAY
            para.space_after = Pt(5)

    if caption:
        cap = slide.shapes.add_textbox(Inches(0.35), Inches(6.85), Inches(12.5), Inches(0.45))
        cp = cap.text_frame.paragraphs[0]
        cp.text = caption
        cp.font.size = Pt(10)
        cp.font.italic = True
        cp.font.color.rgb = GRAY


def add_workflow_slide(prs, title, phases):
    """phases: list of (phase_name, [steps])"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, prs, title)

    col_w = 12.0 / len(phases)
    for i, (phase, steps) in enumerate(phases):
        x = Inches(0.4 + i * col_w)
        w = Inches(col_w - 0.15)

        hdr = slide.shapes.add_shape(1, x, Inches(1.1), w, Inches(0.45))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = GOLD
        hdr.line.fill.background()
        ht = slide.shapes.add_textbox(x + Inches(0.05), Inches(1.15), w - Inches(0.1), Inches(0.35))
        hp = ht.text_frame.paragraphs[0]
        hp.text = phase
        hp.font.bold = True
        hp.font.size = Pt(11)
        hp.font.color.rgb = WHITE
        hp.alignment = PP_ALIGN.CENTER

        body = slide.shapes.add_textbox(x, Inches(1.65), w, Inches(5.2))
        btf = body.text_frame
        btf.word_wrap = True
        for j, step in enumerate(steps):
            para = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
            para.text = f"{j + 1}. {step}"
            para.font.size = Pt(10)
            para.font.color.rgb = GRAY
            para.space_after = Pt(4)


def build_presentation():
    if not SCREENSHOTS.exists():
        SCREENSHOTS.mkdir(parents=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── TITLE ──
    add_title_slide(
        prs,
        "Shree Hari Jewels",
        "Your Online Store & ERP — Complete Step-by-Step Guide\n"
        f"Store: {STORE_URL}",
    )

    add_steps_slide(
        prs,
        "How Your System Works",
        [
            "Your ERP holds all inventory — every product, price, and stock count.",
            "You choose which products to publish to your website (Online Store section).",
            "Customers browse your website — no login needed — and place orders as guests.",
            "When a customer checks out, the ERP automatically reserves stock and creates a Web Order.",
            "You manage and fulfill web orders from ERP → Online Store → Web Orders.",
            "Payment is arranged offline (bank transfer / cash on delivery) — you update payment status in ERP.",
        ],
        note="This guide covers only your website and how it connects to your ERP.",
    )

    # ── PART 1: SET UP YOUR ONLINE STORE (ERP) ──
    add_section_slide(prs, "Part 1", "Set Up Your Online Store in the ERP")

    add_steps_slide(
        prs,
        "Go Live — 5 Steps in the ERP",
        [
            f"Log in to your ERP at {ERP_URL} → Online Store → Store Settings",
            "Enable the online store and set your branding (logo, hero banner, colours, contact info)",
            "Go to Publish Products — toggle products from your inventory to Published",
            "Optionally create Collections to group products (e.g. Wedding, Rings, Necklaces)",
            f"Share your store URL with customers: {STORE_URL}",
        ],
    )

    add_screenshot_slide(
        prs,
        "Step 1 — Online Store Dashboard",
        "09-erp-storefront-dashboard.png",
        steps=[
            "Login to ERP → sidebar: Online Store → Store dashboard",
            "See stats: published products, collections, pending orders",
            "Click View Store to open your live website",
            "Quick links to Settings, Publish Products, Collections, Web Orders",
        ],
        caption="ERP → Online Store → Store dashboard",
    )

    add_screenshot_slide(
        prs,
        "Step 2 — Store Settings (Branding & Enable)",
        "10-erp-store-settings.png",
        steps=[
            "Check Enable online store to make the website public",
            "Set tagline, hero title, hero subtitle, about text",
            "Choose primary & accent colours for your website theme",
            "Add logo URL, banner URL, contact phone, WhatsApp, address",
            "Click Save — changes appear on your website immediately",
        ],
        caption="ERP → Online Store → Store settings",
    )

    add_screenshot_slide(
        prs,
        "Step 3 — Publish Products to Website",
        "11-erp-publish-products.png",
        steps=[
            "Lists all products from your ERP inventory",
            "Only Published products appear on the website",
            "Check stock count — website shows live available stock",
            "Prices update automatically from current gold/silver market rates",
            "Toggle a product to Published → it goes live on the website",
        ],
        caption="ERP → Online Store → Publish products",
    )

    add_screenshot_slide(
        prs,
        "Step 4 — Create Collections (Optional)",
        "12-erp-collections.png",
        steps=[
            "Group products into curated collections (e.g. Featured, Wedding)",
            "Enter collection name and description → Create Collection",
            "Add published products to each collection",
            "Collections appear on website homepage and /collections page",
        ],
        caption="ERP → Online Store → Collections",
    )

    # ── PART 2: CUSTOMER WEBSITE JOURNEY ──
    add_section_slide(prs, "Part 2", "What Your Customer Sees on the Website")

    add_steps_slide(
        prs,
        "Customer Journey — 7 Steps",
        [
            "Customer opens your store URL (or your custom domain)",
            "Browses homepage, product catalog, or collections",
            "Opens a product → sees metal, purity, weight, live price, stock status",
            "Clicks Add to Cart (cart saved in browser — no account needed)",
            "Reviews cart → Proceed to Checkout",
            "Fills guest form: name, mobile, delivery address",
            "Clicks Place Order → sees confirmation with order number",
        ],
        note="No customer login or registration — checkout is fully guest-based.",
    )

    add_screenshot_slide(
        prs,
        "Step 1 — Website Homepage",
        "01-shop-home.png",
        steps=[
            "Hero banner with your tagline and Shop Now button",
            "Featured products pulled from your published inventory",
            "Navigation: Home, Shop, Collections, About",
            "Cart icon shows item count when customer adds products",
        ],
        caption=f"Customer visits: {STORE_URL}",
    )

    add_screenshot_slide(
        prs,
        "Step 2 — Product Catalog (Shop All)",
        "02-shop-products.png",
        steps=[
            "All published products shown in a grid",
            "Customer can search by name",
            "Filter by category, sort by price or name",
            "Only products you marked Published in ERP appear here",
        ],
        caption="Website → Shop → All Products",
    )

    add_screenshot_slide(
        prs,
        "Step 3 — Product Detail Page",
        "03-product-detail.png",
        steps=[
            "Product image, name, SKU, live price",
            "Specs: Metal, Purity, Weight, Availability",
            "Customer selects quantity",
            "Clicks Add to Cart → item saved to browser cart",
        ],
        caption="Website → Product detail → Add to Cart",
    )

    add_screenshot_slide(
        prs,
        "Step 4 — Shopping Cart",
        "04-cart.png",
        steps=[
            "Review items, adjust quantities, or remove items",
            "Subtotal shown — prices from live market rates",
            "Cart persists even if customer closes the browser",
            "Click Proceed to Checkout to continue",
        ],
        caption="Website → Cart",
    )

    add_screenshot_slide(
        prs,
        "Step 5 — Guest Checkout",
        "05-checkout.png",
        steps=[
            "Customer fills: Full name, Mobile (10 digits), Address, City, State, Pincode",
            "Email and order notes are optional",
            "Order summary shows items and total on the right",
            "Payment note: arranged with store after order (no online payment)",
            "Click Place Order to submit",
        ],
        caption="Website → Checkout (guest — no login required)",
    )

    add_screenshot_slide(
        prs,
        "Step 6 — Order Confirmation",
        "08-order-confirmation.png",
        steps=[
            "Thank you page with order number (e.g. WEB-2026-0002)",
            "Shows items, total, and delivery address",
            "Message: store will contact customer for payment & delivery",
            "Customer can continue shopping",
        ],
        caption="Website → Order confirmation",
    )

    add_screenshot_slide(
        prs,
        "Other Website Pages",
        "07-about.png",
        steps=[
            "About page — your business story and contact details",
            "Collections page — curated product groups you create in ERP",
            "All pages use your branding colours and logo from Store Settings",
        ],
        caption="Website → About / Collections",
    )

    # ── PART 3: WHAT HAPPENS IN ERP WHEN ORDER IS PLACED ──
    add_section_slide(prs, "Part 3", "What Happens in Your ERP When a Customer Orders")

    add_steps_slide(
        prs,
        "Backend Flow — When Customer Clicks Place Order",
        [
            "System validates cart items and checks stock is available",
            "Finds or creates a Customer record using the mobile number",
            "Reserves inventory units (status: Reserved — not sold yet)",
            "Creates a WebOrder with order number (WEB-2026-XXXX)",
            "Creates a linked ERP Order for your fulfillment team",
            "Product stock counts update automatically",
            "You see the new order in ERP → Online Store → Web Orders",
        ],
    )

    add_screenshot_slide(
        prs,
        "Step 7 — Manage Web Orders in ERP",
        "13-erp-web-orders.png",
        steps=[
            "New web orders appear here automatically",
            "See customer name, mobile, items, total, date",
            "Update Order Status: Pending → Confirmed → Processing → Shipped → Delivered",
            "Update Payment Status: Unpaid → Paid (after bank transfer / COD)",
            "Contact customer at their mobile to confirm payment & delivery",
        ],
        caption="ERP → Online Store → Web orders",
    )

    add_steps_slide(
        prs,
        "Fulfill a Web Order — Step by Step",
        [
            "New order arrives in Web Orders (status: Pending, payment: Unpaid)",
            "Call/message customer to confirm order and arrange payment",
            "Update status to Confirmed once customer confirms",
            "Collect payment (bank transfer / cash on delivery) → set Payment to Paid",
            "Pack and dispatch the item → update status to Shipped",
            "On delivery → update status to Delivered",
            "Inventory units move from Reserved → Sold when order is fulfilled",
        ],
    )

    add_workflow_slide(
        prs,
        "Complete End-to-End Workflow",
        [
            (
                "YOU (ERP Setup)",
                [
                    "Add stock to inventory",
                    "Verify entry & set prices",
                    "Enable online store",
                    "Publish products",
                    "Share store URL",
                ],
            ),
            (
                "CUSTOMER (Website)",
                [
                    "Browse catalog",
                    "Add to cart",
                    "Guest checkout",
                    "Place order",
                    "Get confirmation",
                ],
            ),
            (
                "YOU (ERP Fulfillment)",
                [
                    "See Web Order",
                    "Confirm with customer",
                    "Collect payment",
                    "Ship item",
                    "Mark Delivered",
                ],
            ),
        ],
    )

    add_steps_slide(
        prs,
        "Website vs In-Store Sale — Key Differences",
        [
            "Website: guest checkout, no login — In-store POS (/sales): staff scans item codes",
            "Website: payment arranged later — In-store: Cash / UPI QR / Card at counter",
            "Website: reserves stock (Reserved) — In-store: marks Sold immediately",
            "Website: creates WebOrder — In-store: creates Sale + GST Invoice",
            "Both use the same inventory and customer records in your ERP",
        ],
        note="Your website and in-store counter share the same inventory — stock is always in sync.",
    )

    add_title_slide(
        prs,
        "You're All Set!",
        f"Your Store: {STORE_URL}\n"
        f"Your ERP: {ERP_URL}\n\n"
        "Manage everything from ERP → Online Store",
    )

    prs.save(str(OUTPUT))
    print(f"Presentation saved to: {OUTPUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
